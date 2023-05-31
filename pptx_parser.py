from pptx import Presentation


class PptxParser:
    """
    Parsing a pptx file and yield the slide's text and index number.
    """
    @classmethod
    def parse(cls, filepath):
        """
        Receives a path to a pptx file and parse the presentation.
        :param filepath: str path to a pptx file
        :yield: slide's text and index number
        """
        prs = Presentation(filepath)
        for index, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                # saves the text from each paragraph in the slide
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run_text = run.text.strip()
                        if run_text:
                            slide_text.append(run_text)

            slide_text = " ".join(slide_text)
            yield index, slide_text

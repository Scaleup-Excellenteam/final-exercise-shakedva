import os
import asyncio
import openai
import json
import time
from pptx import Presentation

FILE_PATH = r'C:\Users\shake\Desktop\College\4th Year\Semester B\Excellenteam\python\Ex\Tests.pptx'
PROMPT = "Please summarize the key points from slide number {slide_index} titled '{file_title}' simply: {slide_text}"
openai.api_key = os.getenv('OPENAI_API_KEY')
BASE_FILE_NAME = os.path.basename(FILE_PATH)


async def parse_slide(slide, slide_index, gpt_outputs):
    slide_text = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run_text = run.text.strip()
                if run_text:
                    slide_text.append(run_text)

    slide_text = " ".join(slide_text)
    gpt_result = await get_model_response(slide_text, slide_index)
    gpt_outputs.append(gpt_result)


async def get_model_response(slide_text, slide_index):
    while True:
        try:
            completion = openai.ChatCompletion.create(
                model="gpt-3.5-turbo",
                messages=[
                    {
                        "role": "user",
                        "content": PROMPT.format(
                            slide_index=slide_index,
                            file_title=BASE_FILE_NAME,
                            slide_text=slide_text
                        )
                    }
                ]
            )
            return completion.choices[0].message.content
        except openai.error.RateLimitError:
            time.sleep(60)


async def main():
    tasks = []
    gpt_outputs = []
    prs = Presentation(FILE_PATH)
    for index, slide in enumerate(prs.slides):
        tasks.append(parse_slide(slide, index, gpt_outputs))
    await asyncio.gather(*tasks)
    with open(f'{os.path.splitext(BASE_FILE_NAME)[0]}_explained.json', 'w') as f:
        f.write(json.dumps(gpt_outputs))


if __name__ == "__main__":
    asyncio.run(main())


